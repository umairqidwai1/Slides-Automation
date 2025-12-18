from pptx import Presentation
import json
import os
from pptx.util import Pt

def extract_original_formatting(shape):
    """Extracts font properties (color, size, bold) from the first run of a shape."""
    if not shape.has_text_frame or len(shape.text_frame.paragraphs) == 0:
        return None
    first_paragraph = shape.text_frame.paragraphs[0]
    if len(first_paragraph.runs) == 0:
        return None
    first_run = first_paragraph.runs[0]
    return {
        "color": first_run.font.color.rgb,
        "size": first_run.font.size,
        "bold": first_run.font.bold,
        "name": first_run.font.name,
    }

def apply_formatting(run, formatting):
    """Applies formatting to a text run."""
    if formatting["color"]:
        run.font.color.rgb = formatting["color"]
    if formatting["size"]:
        run.font.size = formatting["size"]
    run.font.bold = formatting["bold"]
    run.font.name = formatting["name"]

def delete_slide(prs, index):
    """Delete a slide by index (python-pptx workaround)."""
    sldIdLst = prs.slides._sldIdLst
    sldId = list(sldIdLst)[index]
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)

def main():
    prs = Presentation(os.path.join("templates", "Template.pptx"))
    
    with open("content.json", "r", encoding="utf-8") as f:
        slides_data = json.load(f)
    
    # Get original formatting from Slide 2's title and body
    master_slide = prs.slides[1]
    title_formatting = extract_original_formatting(master_slide.shapes[10])  # Title shape
    body_formatting = extract_original_formatting(master_slide.shapes[11])   # Body shape

    # Ensure heading is larger (add 20% to original size)
    if title_formatting and title_formatting["size"]:
        title_formatting["size"] = Pt(title_formatting["size"].pt * 1.2)

    # --- Cover slide (Slide 1) ---
    if slides_data:
        cover = slides_data[0]
        cover_title = cover.get("title", "")
        body_list = cover.get("body", [])
        teacher = body_list[0] if len(body_list) > 0 else ""
        date_str = body_list[1] if len(body_list) > 1 else ""

        slide = prs.slides[0]
        
        # Title (shape[11])
        title_shape = slide.shapes[11].text_frame
        title_shape.clear()
        p = title_shape.paragraphs[0] if title_shape.paragraphs else title_shape.add_paragraph()
        p.text = f"{cover_title}"
        for run in p.runs:
            apply_formatting(run, title_formatting)
        
        # Teacher (shape[10])
        teacher_shape = slide.shapes[10].text_frame
        teacher_shape.clear()
        p = teacher_shape.paragraphs[0] if teacher_shape.paragraphs else teacher_shape.add_paragraph()
        p.text = f"{teacher}"
        for run in p.runs:
            apply_formatting(run, body_formatting)

        # Date (shape[12])
        date_shape = slide.shapes[12].text_frame
        date_shape.clear()
        p = date_shape.paragraphs[0] if date_shape.paragraphs else date_shape.add_paragraph()
        p.text = f"{date_str}"
        for run in p.runs:
            apply_formatting(run, body_formatting)

    # Update existing slides (content starts at JSON index 1)
    for idx, slide_info in enumerate(slides_data[1:]):
        slide = prs.slides[idx + 1]  # Target Slide 2 and onward
        
        # Update Title (index 10)
        title_shape = slide.shapes[10].text_frame
        title_shape.clear()
        p = title_shape.paragraphs[0] if title_shape.paragraphs else title_shape.add_paragraph()
        p.text = slide_info["title"]
        for run in p.runs:
            apply_formatting(run, title_formatting)
        
        # Update Body (index 11)
        body_shape = slide.shapes[11].text_frame
        body_shape.clear()
        for line in slide_info["body"]:
            p = body_shape.add_paragraph()
            p.text = line.strip()
            p.level = 1 if line.startswith("-") else 0
            for run in p.runs:
                apply_formatting(run, body_formatting)

    # --- Trim any extra slides beyond JSON count ---
    target_count = len(slides_data)
    while len(prs.slides) > target_count + 1:
        delete_slide(prs, len(prs.slides) - 2)

    prs.save("Generated_Presentation.pptx")
    print("Done! Presentation saved.")

if __name__ == "__main__":
    main()
